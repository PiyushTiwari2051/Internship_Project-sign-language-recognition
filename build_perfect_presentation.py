import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_perfect_presentation():
    template_path = r'c:\Users\HP\Downloads\Internship Presentation PPTs Template Final.pptx'
    output_path = r'c:\Users\HP\Downloads\Piyush_Tiwari_Sign_Language_Recognition_Final_PPT.pptx'
    workspace_copy = r'Piyush_Tiwari_Sign_Language_Recognition_Final_PPT.pptx'

    prs = Presentation(template_path)

    # Color Palette
    NAVY = RGBColor(15, 23, 42)          # #0f172a
    BLUE = RGBColor(13, 110, 253)        # #0d6efd
    DARK_GRAY = RGBColor(51, 65, 85)     # #334155
    LIGHT_BLUE = RGBColor(224, 242, 254) # #e0f2fe
    ACCENT_PURPLE = RGBColor(139, 92, 246) # #8b5cf6
    GREEN = RGBColor(16, 185, 129)       # #10b981
    WHITE = RGBColor(255, 255, 255)
    RED = RGBColor(239, 68, 68)          # #ef4444

    # ==================== SLIDE 1: Modify In-Place ====================
    slide1 = prs.slides[0]
    
    # Update Project Name in Shape 7
    shape7 = slide1.shapes[7]
    tf7 = shape7.text_frame
    for p in tf7.paragraphs:
        if 'Project Name' in p.text:
            p.text = 'Real-Time Sign Language Recognition & Translation'
            p.font.name = 'Calibri'
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 0, 0) # RED as in template

    # Update Student Name & Roll No in Shape 1
    shape1 = slide1.shapes[1]
    tf1 = shape1.text_frame
    full_txt = tf1.text
    full_txt = full_txt.replace('Name of the Student', 'Piyush Tiwari')
    full_txt = full_txt.replace('Roll no.', '2100270100000')
    tf1.text = full_txt
    for p in tf1.paragraphs:
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = NAVY

    # Delete existing template slides 2 to end while keeping Slide 1
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1]

    def add_slide_header(slide, title_text, slide_num):
        # Slide Title 32pt Bold
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Calibri'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = NAVY

        # Slide Number Footer
        num_box = slide.shapes.add_textbox(Inches(8.5), Inches(6.8), Inches(1.0), Inches(0.4))
        num_tf = num_box.text_frame
        p_num = num_tf.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.name = 'Calibri'
        p_num.font.size = Pt(14)
        p_num.font.color.rgb = DARK_GRAY
        p_num.alignment = PP_ALIGN.RIGHT

    def add_bullets(slide, items, left=0.6, top=1.3, width=8.8, height=5.3):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        
        for idx, item in enumerate(items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.space_after = Pt(6) # Controlled spacing to prevent overflow
            
            if isinstance(item, tuple):
                heading, body = item
                p.text = "• " + heading + ": "
                p.font.name = 'Calibri'
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = NAVY
                
                run = p.add_run()
                run.text = body
                run.font.name = 'Calibri'
                run.font.size = Pt(18)
                run.font.bold = False
                run.font.color.rgb = DARK_GRAY
            else:
                p.text = "• " + item
                p.font.name = 'Calibri'
                p.font.size = Pt(20)
                p.font.bold = False
                p.font.color.rgb = DARK_GRAY

    # ==================== SLIDE 2: Outline & Contents ====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide2, "Outline & Table of Contents", 2)
    outline_items = [
        ("1. Introduction & Objectives", "Domain of work, assistive AI technology & project goals"),
        ("2. Problem Statement", "Challenges in sign language translation & pixel-level limitations"),
        ("3. Requirement Analysis", "Hardware setup & software stack (MediaPipe, TensorFlow, OpenCV)"),
        ("4. System Architecture", "End-to-end data pipeline flow & system module design"),
        ("5. Keypoint Extraction & Normalization", "Spatial landmark extraction (Hands & Pose)"),
        ("6. Dataset & Feature Engineering", "Sliding window aggregation (18 statistical features)"),
        ("7. Sequence LSTM Classifier", "Multi-layer LSTM model architecture & loss optimization"),
        ("8. Desktop GUI & Audio Translation", "Native Python Tkinter application & Pyttsx3 speech synthesis"),
        ("9. Results & Performance Analysis", "Accuracy (98.67%), F1-Score (0.99), & per-class precision table"),
        ("10. Conclusion & References", "Summary of achievements, future scope & academic references")
    ]
    add_bullets(slide2, outline_items)

    # ==================== SLIDE 3: Introduction & Objectives ====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide3, "Introduction & Objectives", 3)
    intro_items = [
        ("Domain of Work", "Computer Vision, Deep Learning & Accessibility Assistive Technology"),
        ("Social Relevance", "Over 70 million deaf individuals globally rely on sign language for daily communication"),
        ("Primary Objective", "Develop a real-time, non-invasive translation system converting webcam gestures into text & voice"),
        ("Key Approach", "Extract compact 3D spatial keypoints per frame and classify keypoint sequences using an LSTM model"),
        ("Sign-to-Speech Integration", "Integrate Text-to-Speech (TTS) audio output for seamless real-time communication")
    ]
    add_bullets(slide3, intro_items)

    # ==================== SLIDE 4: Problem Statement & Motivation ====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide4, "Problem Statement & Motivation", 4)
    prob_items = [
        ("Raw Video Heavy & Noisy", "Raw video frames contain high spatial redundancy, background noise, and heavy memory cost"),
        ("Limitations of Existing CNN Methods", "Direct frame-by-frame CNN classifiers fail on temporal gesture dynamics and suffer from high latency"),
        ("Signer Position Dependency", "Raw pixel models degrade when signer distance, clothing, or background lighting changes"),
        ("Our Proposed Solution", "Replace raw pixels with 3D pose/hand keypoint sequences, ensuring position scale invariance"),
        ("Real-Time Constraint", "Achieve fast sub-50ms inference latency suitable for continuous live webcam translation")
    ]
    add_bullets(slide4, prob_items)

    # ==================== SLIDE 5: Requirement Analysis ====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide5, "Requirement Analysis (S/W & H/W)", 5)
    req_items = [
        ("Hardware Requirements", "Standard HD Webcam (640x480 resolution @ 15+ FPS), Intel Core i5/i7, 8GB RAM"),
        ("Programming Language", "Python 3.11 (Core logic, OpenCV capture, & deep learning execution)"),
        ("Computer Vision Library", "Google MediaPipe 0.10.21 (Fingertip & upper-body pose landmark extraction)"),
        ("Deep Learning Framework", "TensorFlow 2.17 & Keras 3.3 (Multi-layer LSTM model training & inference)"),
        ("User Interface & Audio Engine", "Python Native Tkinter GUI framework & Pyttsx3 Text-to-Speech audio engine")
    ]
    add_bullets(slide5, req_items)

    # ==================== SLIDE 6: Interactive Flowchart ====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide6, "System Methodology & Data Flowchart", 6)

    # Professional Flowchart Grid (2 rows x 3 cols + Output box)
    flow_steps = [
        ("STAGE 1", "Webcam Capture\n(640x480 @ 15 FPS)", BLUE),
        ("STAGE 2", "MediaPipe Extraction\n(Hands & Pose Joints)", ACCENT_PURPLE),
        ("STAGE 3", "Sliding Window\n(0.5s Feature Aggregation)", BLUE),
        ("STAGE 4", "Normalization\n(StandardScaler Normalizer)", NAVY),
        ("STAGE 5", "LSTM Sequence Net\n(5-Window Input Tensor)", ACCENT_PURPLE),
        ("STAGE 6", "GUI & Audio Speech\n(Hero Box + Pyttsx3 TTS)", GREEN)
    ]

    left_base = 0.6
    top_base = 1.5
    w_box = 2.4
    h_box = 1.8

    for idx, (st_num, st_name, col) in enumerate(flow_steps):
        r = idx // 3
        c = idx % 3
        
        x = left_base + c * (w_box + 0.6)
        y = top_base + r * (h_box + 0.6)

        # Draw card container
        shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w_box), Inches(h_box))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE
        shape.line.color.rgb = col
        shape.line.width = Pt(2.5)

        tf = shape.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = st_num
        p0.font.name = 'Calibri'
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = col
        p0.alignment = PP_ALIGN.CENTER

        p1 = tf.add_paragraph()
        p1.text = st_name
        p1.font.name = 'Calibri'
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.alignment = PP_ALIGN.CENTER

        # Draw Arrow connector between cols
        if c < 2:
            arrow_x = x + w_box + 0.05
            arrow_y = y + h_box / 2.0 - 0.15
            arr = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(arrow_y), Inches(0.5), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = col
            arr.line.fill.background()

    # Flowchart Caption
    cap6 = slide6.shapes.add_textbox(Inches(0.6), Inches(6.1), Inches(8.8), Inches(0.5))
    p_cap6 = cap6.text_frame.paragraphs[0]
    p_cap6.text = "Figure 1: End-to-End Real-Time Gesture Recognition Data Flow Pipeline"
    p_cap6.font.name = 'Calibri'
    p_cap6.font.size = Pt(16)
    p_cap6.font.bold = True
    p_cap6.font.color.rgb = DARK_GRAY
    p_cap6.alignment = PP_ALIGN.CENTER

    # ==================== SLIDE 7: Keypoint Extraction & Normalization ====================
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide7, "Keypoint Extraction & Normalization", 7)
    key_items = [
        ("Hand Landmarks Extraction", "Extracts 5 fingertip joints (thumb, index, middle, ring, pinky) per hand (X, Y, Z coordinates)"),
        ("Upper-Body Pose Landmarks", "Extracts 6 upper-body joints (shoulders, elbows, wrists) for posture tracking"),
        ("Statistical Feature Aggregation", "Computes Mean and Standard Deviation per 0.5s sliding window (~7-8 frames)"),
        ("Scale & Position Invariance", "StandardScaler normalization ensures robustness regardless of signer camera distance"),
        ("Compact Vector Dimension", "Reduces full video frame data into a compact 18-element feature vector per window")
    ]
    add_bullets(slide7, key_items)

    # ==================== SLIDE 8: Dataset & Feature Engineering ====================
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide8, "Dataset & Feature Details", 8)
    ds_items = [
        ("Dataset Size & Support", "Curated dataset containing 3,690 gesture sequence samples across 11 sign classes"),
        ("Target Sign Vocabulary", "Afternoon, Age, Boy, Country, Day, Monday, Name, Night, People, Person, Time"),
        ("Temporal Window Length", "Fixed 0.5-second sliding window duration (target sampling rate: 15 FPS)"),
        ("Sequence Stacking", "5 consecutive feature windows stacked into an input matrix of shape (1, 5, 18)"),
        ("Database Storage", "Stored using SQLite schema tracking video metadata, frame landmarks, and feature matrices")
    ]
    add_bullets(slide8, ds_items)

    # ==================== SLIDE 9: Sequence LSTM Architecture ====================
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide9, "Sequence LSTM Model Architecture", 9)
    model_items = [
        ("Input Layer Shape", "(5, 18) representing 5 time-step windows of 18 normalized keypoint features"),
        ("LSTM Layer 1", "128 Recurrent LSTM Units with Dropout (0.3) & Batch Normalization for regularization"),
        ("LSTM Layer 2", "64 Recurrent LSTM Units capturing higher-level temporal gesture dynamics"),
        ("Dense Classification Layer", "64 Dense Neurons with ReLU activation & Dropout (0.3)"),
        ("Output Softmax Layer", "11 Output Neurons producing Softmax confidence probabilities over sign vocabulary")
    ]
    add_bullets(slide9, model_items)

    # ==================== SLIDE 10: Desktop GUI Implementation ====================
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide10, "Desktop GUI Implementation (Tkinter)", 10)
    gui_items = [
        ("100% Native Python Desktop GUI", "Built using Python Tkinter & Pillow (cv2main.py) for a clean, natural software feel"),
        ("Live Video Viewport", "Displays real-time OpenCV camera feed with MediaPipe landmark skeleton overlays"),
        ("Prediction Hero Display", "Prominently displays top predicted gesture text in a clear, bold hero card"),
        ("Confidence Progress Meters", "Native ttk.Progressbar widgets detailing real-time top-3 gesture probabilities"),
        ("Dataset Management Tabs", "Includes dedicated tabs for Recording Videos, Importing Files, and Dataset Video Gallery")
    ]
    add_bullets(slide10, gui_items)

    # ==================== SLIDE 11: Active Hand Gating & Sign-to-Speech ====================
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide11, "Active Hand Gating & Sign-to-Speech", 11)
    gate_items = [
        ("Active Hand Detection Gating", "Prevents false predictions when resting; displays [ NO HAND DETECTED ] when inactive"),
        ("Confidence Thresholding", "Requires a minimum 45% probability threshold before displaying sign gesture output"),
        ("Text-to-Speech (TTS) Integration", "Asynchronous Pyttsx3 speech engine speaks recognized gesture words aloud"),
        ("Audio Toggle Control", "Includes a GUI checkbox enabling/disabling voice output for user preference"),
        ("Sub-30ms Inference Latency", "Ensures smooth live translation without lag or freezing during continuous signing")
    ]
    add_bullets(slide11, gate_items)

    # ==================== SLIDE 12: Empirical Results & Model Evaluation ====================
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide12, "Empirical Results & Model Evaluation", 12)
    eval_items = [
        ("Overall Test Accuracy", "Achieved 98.67% (0.9867) classification accuracy on 3,690 test evaluation samples"),
        ("Overall Test Loss", "Achieved low cross-entropy loss of 0.0459 during test set evaluation"),
        ("Macro Average F1-Score", "Macro F1-Score of 0.99 demonstrating balanced per-class model accuracy"),
        ("Weighted Average F1-Score", "Weighted F1-Score of 0.99 confirming strong generalizability across all gesture classes"),
        ("Inference Latency", "Average frame processing time under 30ms, operating smoothly at 15+ FPS")
    ]
    add_bullets(slide12, eval_items)

    # ==================== SLIDE 13: Per-Class Evaluation Table ====================
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide13, "Per-Class Evaluation Methodology Table", 13)

    rows, cols = 13, 5
    left, top, width, height = Inches(0.6), Inches(1.3), Inches(8.8), Inches(4.7)
    table_shape = slide13.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table_data = [
        ["Sign Gesture", "Precision", "Recall", "F1-Score", "Test Samples"],
        ["Afternoon", "0.99", "1.00", "0.99", "266"],
        ["Age", "0.98", "0.99", "0.99", "387"],
        ["Boy", "1.00", "1.00", "1.00", "290"],
        ["Country", "0.98", "0.95", "0.96", "428"],
        ["Day", "0.99", "1.00", "0.99", "427"],
        ["Monday", "1.00", "1.00", "1.00", "212"],
        ["Name", "0.99", "0.98", "0.98", "401"],
        ["Night", "1.00", "1.00", "1.00", "322"],
        ["People", "1.00", "1.00", "1.00", "422"],
        ["Person", "1.00", "1.00", "1.00", "238"],
        ["Time", "0.93", "0.96", "0.95", "297"],
        ["OVERALL AVG", "0.99", "0.99", "0.99", "3,690"]
    ]

    for r_idx, row in enumerate(table_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Calibri'
            p.font.size = Pt(13) if r_idx > 0 else Pt(14)
            p.font.bold = (r_idx == 0 or r_idx == 12)
            p.alignment = PP_ALIGN.CENTER
            if r_idx == 0:
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif r_idx == 12:
                p.font.color.rgb = NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            else:
                p.font.color.rgb = DARK_GRAY

    cap13 = slide13.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(8.8), Inches(0.4))
    p_c13 = cap13.text_frame.paragraphs[0]
    p_c13.text = "Table 1: Per-Class Precision, Recall, F1-Score, and Support Sample Counts"
    p_c13.font.name = 'Calibri'
    p_c13.font.size = Pt(15)
    p_c13.font.bold = True
    p_c13.font.color.rgb = DARK_GRAY
    p_c13.alignment = PP_ALIGN.CENTER

    # ==================== SLIDE 14: Observations & Result Analysis ====================
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide14, "Observations & Result Analysis", 14)
    obs_items = [
        ("Perfect Class Differentiation", "Gestures like 'Boy', 'Monday', 'Night', 'People', and 'Person' achieved 1.00 F1-score"),
        ("Minor Gesture Overlap", "Slight overlap observed between 'Country' (0.96 F1) and 'Time' (0.95 F1) due to similar arm positions"),
        ("Signer Position Independence", "Keypoint spatial normalization successfully eliminated distance & clothing variations"),
        ("Resting State Stability", "Hand gating prevented false random predictions when signer arms were at rest"),
        ("Audio Responsiveness", "Text-to-speech audio feedback triggered smoothly within ~0.4s of gesture completion")
    ]
    add_bullets(slide14, obs_items)

    # ==================== SLIDE 15: Conclusion & Summary ====================
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide15, "Conclusion & Summary", 15)
    conc_items = [
        ("Successful System Delivery", "Built a complete real-time sign language gesture recognition & audio translation system"),
        ("High Model Performance", "Achieved 98.67% overall test accuracy utilizing MediaPipe keypoints & LSTM architecture"),
        ("Sign-to-Speech Integration", "Successfully integrated live text display and Pyttsx3 speech audio output for accessibility"),
        ("Native Desktop Application", "Developed a 100% native Python Tkinter GUI featuring Live View, Recording, & Gallery tabs"),
        ("Real-World Practicality", "Provides an effective, low-latency communication tool for deaf & hard-of-hearing assistance")
    ]
    add_bullets(slide15, conc_items)

    # ==================== SLIDE 16: Future Scope & Recommendations ====================
    slide16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide16, "Future Scope & Recommendations", 16)
    fut_items = [
        ("Vocabulary Expansion", "Scale vocabulary from 11 signs to full datasets like WLASL (2,000+ signs) or INCLUDE dataset"),
        ("Sentence-Level Translation", "Extend from word-level gestures to full continuous sentence translation using Transformer models"),
        ("Mobile & Edge Deployment", "Export model to TensorFlow Lite (TFLite) for deployment on mobile devices & Raspberry Pi"),
        ("Multi-Language Speech Output", "Add multi-language TTS translation options (e.g. Hindi, Spanish, French audio output)"),
        ("Dataset Mirror Augmentation", "Incorporate mirror data augmentation for left/right hand signer generalizability")
    ]
    add_bullets(slide16, fut_items)

    # ==================== SLIDE 17: References & Bibliography ====================
    slide17 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide17, "References & Bibliography", 17)
    ref_items = [
        "1. Lugaresi et al., 'MediaPipe: A Framework for Building Perception Pipelines', arXiv:1906.08172, 2019.",
        "2. Hochreiter, S., & Schmidhuber, J., 'Long Short-Term Memory', Neural Computation, 9(8), 1735-1780, 1997.",
        "3. Li, D., et al., 'Word-Level Deep Sign Language Recognition', IEEE/CVF Conference on Computer Vision (CVPR), 2020.",
        "4. Metehan Ozdeniz, 'Sign Language Recognition Dataset & Benchmark', Kaggle Dataset Repository, 2024.",
        "5. Google Keras 3.3 & TensorFlow 2.17 Official Documentation and API Reference, 2026.",
        "6. OpenCV Open Source Computer Vision Library Documentation, 2026."
    ]
    add_bullets(slide17, ref_items)

    # ==================== SLIDE 18: Acknowledgement & Q/A ====================
    slide18 = prs.slides.add_slide(prs.slide_layouts[6])
    
    tb18 = slide18.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(4.5))
    tf18 = tb18.text_frame
    tf18.word_wrap = True

    p18_1 = tf18.paragraphs[0]
    p18_1.text = "Thank You!"
    p18_1.font.name = 'Calibri'
    p18_1.font.size = Pt(44)
    p18_1.font.bold = True
    p18_1.font.color.rgb = NAVY
    p18_1.alignment = PP_ALIGN.CENTER

    p18_2 = tf18.add_paragraph()
    p18_2.text = "Questions & Discussion"
    p18_2.font.name = 'Calibri'
    p18_2.font.size = Pt(28)
    p18_2.font.bold = True
    p18_2.font.color.rgb = BLUE
    p18_2.alignment = PP_ALIGN.CENTER

    p18_3 = tf18.add_paragraph()
    p18_3.text = "\nPresented By: Piyush Tiwari\nRoll No: 2100270100000 | 7th Sem, CSE-1\nIMS Engineering College, Ghaziabad"
    p18_3.font.name = 'Calibri'
    p18_3.font.size = Pt(20)
    p18_3.font.bold = False
    p18_3.font.color.rgb = DARK_GRAY
    p18_3.alignment = PP_ALIGN.CENTER

    # Save Presentation
    prs.save(output_path)
    prs.save(workspace_copy)
    print(f"PERFECT presentation saved successfully to:\n1. {output_path}\n2. {workspace_copy}")

if __name__ == '__main__':
    build_perfect_presentation()
